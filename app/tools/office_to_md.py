from pathlib import Path
from typing import Any


class OfficeConversionError(Exception):
    """Office 文档转换失败。"""


def convert_docx_to_markdown(file_path: Path, media_dir: Path | None = None) -> tuple[str, list[dict[str, Any]]]:
    """将 docx 文档转换为 Markdown。

    Args:
        file_path: 待转换的 docx 文件路径。
        media_dir: 可选媒体输出目录;当前阶段仅保留参数以便后续扩展图片抽取。

    Returns:
        Markdown 文本和媒体元数据列表。

    Raises:
        OfficeConversionError: 当依赖缺失或转换失败时抛出。
    """
    try:
        import mammoth
    except ImportError as exc:
        raise OfficeConversionError("缺少 mammoth 依赖,无法解析 docx 附件。") from exc
    try:
        with file_path.open("rb") as file_obj:
            result = mammoth.convert_to_markdown(file_obj)
    except Exception as exc:
        raise OfficeConversionError("docx 附件解析失败。") from exc
    return str(result.value or ""), []


def convert_pptx_to_markdown(file_path: Path, media_dir: Path | None = None) -> tuple[str, list[dict[str, Any]]]:
    """将 pptx 文档转换为 Markdown。

    Args:
        file_path: 待转换的 pptx 文件路径。
        media_dir: 可选媒体输出目录;当前阶段仅保留参数以便后续扩展图片抽取。

    Returns:
        Markdown 文本和媒体元数据列表。

    Raises:
        OfficeConversionError: 当依赖缺失或转换失败时抛出。
    """
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise OfficeConversionError("缺少 python-pptx 依赖,无法解析 pptx 附件。") from exc
    try:
        presentation = Presentation(str(file_path))
    except Exception as exc:
        raise OfficeConversionError("pptx 附件解析失败。") from exc

    sections: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        sections.append(f"## 第 {index} 页")
        for shape in slide.shapes:
            if getattr(shape, "has_table", False):
                sections.append(_table_to_markdown(shape.table))
            elif getattr(shape, "has_text_frame", False):
                text = "\n".join(paragraph.text for paragraph in shape.text_frame.paragraphs if paragraph.text.strip())
                if text.strip():
                    sections.append(text.strip())
        notes_slide = getattr(slide, "notes_slide", None)
        notes_text_frame = getattr(getattr(notes_slide, "notes_text_frame", None), "text", "") if notes_slide else ""
        if notes_text_frame and notes_text_frame.strip():
            sections.append(f"备注：\n{notes_text_frame.strip()}")
    return "\n\n".join(sections), []


def _table_to_markdown(table: Any) -> str:
    """将 pptx 表格转换为简单 Markdown 表格。"""
    rows = [[cell.text.strip().replace("\n", " ") for cell in row.cells] for row in table.rows]
    if not rows:
        return ""
    header = rows[0]
    separator = ["---"] * len(header)
    lines = ["| " + " | ".join(header) + " |", "| " + " | ".join(separator) + " |"]
    for row in rows[1:]:
        lines.append("| " + " | ".join(row) + " |")
    return "\n".join(lines)
